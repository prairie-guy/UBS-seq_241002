#!/usr/bin/env python3

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from typing import Union

def add_imgs_to_ppt(in_path: Union[str, Path], out_path: Union[str, Path]) -> None:
    """
    Create a PowerPoint presentation with images from a specified directory.

    This function takes images from the input directory, adds them to slides
    in a PowerPoint presentation, and saves the presentation to the specified
    output file.

    Args:
        in_path (Union[str, Path]): Path to the input directory containing images.
        out_path (Union[str, Path]): Path where the output PowerPoint file will be saved.

    Returns:
        None

    Usage:
        From within Python:
        >>> add_imgs_to_ppt('./images', 'output.pptx')

        From command line:
        $ ./mk_test_pptx.py ./images output.pptx

    Note:
        Supported image formats are: .png, .jpg, .jpeg, .gif, .bmp
        Hidden files (starting with '.') are ignored.
    """
    in_path = Path(in_path)
    out_path = Path(out_path)

    prs = Presentation()
    blank = prs.slide_layouts[6]
    w, h = prs.slide_width, prs.slide_height
    
    for img in sorted(in_path.glob('*.*')):
        if img.name.startswith('.'):
            continue
        if img.suffix.lower() not in ('.png', '.jpg', '.jpeg', '.gif', '.bmp'):
            continue
        slide = prs.slides.add_slide(blank)
        pic = slide.shapes.add_picture(str(img), 0, 0)
        scale = min(w / pic.width, h / pic.height)
        pic.width, pic.height = int(pic.width * scale), int(pic.height * scale)
        pic.left, pic.top = (w - pic.width) // 2, (h - pic.height) // 2
    
    prs.save(str(out_path))
    print(f"Saved {len(prs.slides)} slides to {out_path}")

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print(f"Usage: {sys.argv[0]} <input_directory> <output_pptx_file>")
        print("Example: ./mk_test_pptx.py ./images output.pptx")
        sys.exit(1)
    
    in_path = sys.argv[1]
    out_path = sys.argv[2]
    
    add_imgs_to_ppt(in_path, out_path)